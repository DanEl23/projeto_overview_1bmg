import os
import sys
import json
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Cm
from pptx.enum.shapes import MSO_CONNECTOR
import re  
import docx

    
def create_apresentation(artista, artist, mes_foco):
    prs = Presentation()

    # SLIDE CAPA ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Definir o tamanho do slide para widescreen (16:9)
    prs.slide_width = Inches(13.333)  # 13.333in é o equivalente a 16:9 em widescreen
    prs.slide_height = Inches(7.5)  # Altura proporcional para o tamanho widescreen

    # Caminho para a imagem de fundo do slide de capa
    background_image_path = 'resources/Imagens Template Relatório Mensal/image35.png'

    # Adicionar um slide de capa
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Adicionar a imagem de fundo
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a imagem
    image_path = 'resources/Imagens Template Relatório Mensal/image3.png'
    slide.shapes.add_picture(image_path, Cm(1.42), Cm(1.21), width=Cm(2.89), height=Cm(2.06))

    # Adicionar a primeira caixa de texto (título)
    textbox_title = slide.shapes.add_textbox(Cm(1.04), Cm(11.28), Cm(25.2), Cm(3.68))
    text_frame_title = textbox_title.text_frame
    p_title = text_frame_title.paragraphs[0]
    run_title = p_title.add_run()
    run_title.text = "Monthly Report"

    # Configuração da fonte do título
    font_title = run_title.font
    font_title.size = Pt(56)
    font_title.bold = True
    font_title.name = 'DM Sans'
    font_title.color.rgb = RGBColor(255, 255, 255)  # Cor do título em formato RGB hexadecimal

    # Adicionar a segunda caixa de texto (subtítulo)
    textbox_subtitle = slide.shapes.add_textbox(Cm(1.04), Cm(13.75), Cm(13.31), Cm(1.28))
    text_frame_subtitle = textbox_subtitle.text_frame
    p_subtitle = text_frame_subtitle.paragraphs[0]
    run_subtitle = p_subtitle.add_run()
    run_subtitle.text = mes_foco

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(1.34), Cm(15.61), Cm(17), Cm(15.61))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(255, 255, 255)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Configuração da fonte do subtítulo
    font_subtitle = run_subtitle.font
    font_subtitle.size = Pt(24)
    font_subtitle.bold = True
    font_subtitle.name = 'DM Sans'
    font_subtitle.color.rgb = RGBColor(255, 255, 255)  # Cor do subtítulo em formato RGB hexadecimal (#ffffff)

    # Adicionar a caixa de artista
    textbox_artist = slide.shapes.add_textbox(Cm(1.04), Cm(16.57), Cm(31.98), Cm(1.54))
    text_frame_artist = textbox_artist.text_frame
    p_artista = text_frame_artist.paragraphs[0]
    run_artist = p_artista.add_run()
    run_artist.text = artist

    # Configuração da fonte do subtítulo
    font_artist = run_artist.font
    font_artist.size = Pt(30)
    font_artist.bold = True
    font_artist.name = 'DM Sans'
    font_artist.color.rgb = RGBColor(255, 255, 255)  # Cor do subtítulo em formato RGB hexadecimal (#ffffff)
   
    
    #SLIDE TERMOMETRO ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path_1 = 'resources/Imagens Template Relatório Mensal/image7.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path_1, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a imagem
    logo_path = f'resources/termometro/quadro1.png'
    slide.shapes.add_picture(logo_path, Cm(-0.33), Cm(1.79), width=Cm(2.6), height=Cm(7.58))

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(-2.84), Cm(5.03), Cm(7.58), Cm(1.11))
    textbox_first.rotation = 270
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    p_first.alignment = PP_ALIGN.CENTER
    run_first = p_first.add_run()
    run_first.text = "Channel performance"
    font_first = run_first.font
    font_first.size = Pt(14)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(227, 255, 62)  # Cor do texto em azul escuro

    # Adicionar a imagem
    logo_path = f'resources/termometro/quadro2.png'
    slide.shapes.add_picture(logo_path, Cm(-0.33), Cm(10.71), width=Cm(2.6), height=Cm(7.58))

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(-2.83), Cm(13.94), Cm(7.57), Cm(1.11))
    textbox_first.rotation = 270 
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    p_first.alignment = PP_ALIGN.CENTER
    run_first = p_first.add_run()
    run_first.text = "Main Information"
    font_first = run_first.font
    font_first.size = Pt(14)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(227, 255, 62)  # Cor do texto em azul escuro

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(26.3), Cm(3.23), Cm(7.14), Cm(1.37))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    p_first.alignment = PP_ALIGN.CENTER
    run_first = p_first.add_run()
    run_first.text = "Highlights of the month"
    font_first = run_first.font
    font_first.size = Pt(20)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(25, 46, 104)  # Cor do texto em azul escuro

    # Adicionar a imagem
    logo_path = f'resources/termometro/quadro3.png'
    slide.shapes.add_picture(logo_path, Cm(26.72), Cm(4.43), width=Cm(6.45), height=Cm(1.11))

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(26.71), Cm(4.53), Cm(6.45), Cm(0.85))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    p_first.alignment = PP_ALIGN.CENTER
    run_first = p_first.add_run()
    run_first.text = f"{mes_foco}"
    font_first = run_first.font
    font_first.size = Pt(14)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(25, 46, 104)  # Cor do texto em azul escuro

    # Adicionar a imagem
    logo_path = f'resources/termometro/rostos.png'
    slide.shapes.add_picture(logo_path, Cm(3.92), Cm(2.58), width=Cm(19.13), height=Cm(4.1))

    # Adicionar a imagem
    logo_path = f'resources/termometro/cobertura.png'
    slide.shapes.add_picture(logo_path, Cm(3.82), Cm(2.48), width=Cm(4.28), height=Cm(4.28))

    # Adicionar a imagem
    logo_path = f'resources/termometro/cobertura.png'
    slide.shapes.add_picture(logo_path, Cm(8.82), Cm(2.48), width=Cm(4.28), height=Cm(4.28))

    # Adicionar a imagem
    logo_path = f'resources/termometro/cobertura.png'
    slide.shapes.add_picture(logo_path, Cm(13.83), Cm(2.48), width=Cm(4.28), height=Cm(4.28))

    # Adicionar a imagem
    logo_path = f'resources/termometro/cobertura.png'
    slide.shapes.add_picture(logo_path, Cm(18.91), Cm(2.48), width=Cm(4.28), height=Cm(4.28))


    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(3.71), Cm(7.08), Cm(4.43), Cm(1.11))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    p_first.alignment = PP_ALIGN.CENTER
    run_first = p_first.add_run()
    run_first.text = "In fall"
    font_first = run_first.font
    font_first.size = Pt(14)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(153, 153, 153)  # Cor do texto em azul escuro

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(8.72), Cm(7.08), Cm(4.43), Cm(1.11))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    p_first.alignment = PP_ALIGN.CENTER
    run_first = p_first.add_run()
    run_first.text = "Attention!"
    font_first = run_first.font
    font_first.size = Pt(14)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(153, 153, 153)  # Cor do texto em azul escuro

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(13.74), Cm(7.08), Cm(4.43), Cm(1.11))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    p_first.alignment = PP_ALIGN.CENTER
    run_first = p_first.add_run()
    run_first.text = "Stable"
    font_first = run_first.font
    font_first.size = Pt(14)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(153, 153, 153)  # Cor do texto em azul escuro

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(18.92), Cm(7.08), Cm(4.43), Cm(1.11))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    p_first.alignment = PP_ALIGN.CENTER
    run_first = p_first.add_run()
    run_first.text = "We grew up!"
    font_first = run_first.font
    font_first.size = Pt(14)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(153, 153, 153)  # Cor do texto em azul escuro

    # Adicionar a imagem
    logo_path = f'resources/termometro/linha.png'
    slide.shapes.add_picture(logo_path, Cm(4.29), Cm(8.58), width=Cm(19.06), height=Cm(0.27))

    # Adicionar a imagem
    logo_path = f'resources/termometro/gota.png'
    slide.shapes.add_picture(logo_path, Cm(15.56), Cm(8.14), width=Cm(0.8), height=Cm(1.23))

    # Adicionar a imagem
    logo_path = f'resources/termometro/texto1.png'
    slide.shapes.add_picture(logo_path, Cm(3.31), Cm(10.86), width=Cm(25.3), height=Cm(7.58))

    # Adicionar a imagem
    logo_path = f'resources/termometro/logo.png'
    slide.shapes.add_picture(logo_path, Cm(24.91), Cm(3.55), width=Cm(1.47), height=Cm(1.98))

    # Adicionar a terceira caixa de texto do novo slide
    textbox_third = slide.shapes.add_textbox(Cm(3.71), Cm(11.3), Cm(11.01), Cm(1.64))
    text_frame_third = textbox_third.text_frame
    text_frame_third.word_wrap = True
    text_frame_third.auto_size = True
    p_third = text_frame_third.paragraphs[0]
    run_third = p_third.add_run()
    run_third.text = "Texto" 
    font_third = run_third.font
    font_third.size = Pt(14)
    font_third.name = 'DM Sans'
    font_third.color.rgb = RGBColor(0, 0, 0)  # Cor do texto em preto


    # SLIDE RESUMO ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image15.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Analysis Summary"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(11.32), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))


    # Adicionar a terceira caixa de texto do novo slide
    textbox_fourth = slide.shapes.add_textbox(Cm(0.50), Cm(18.34), Cm(11), Cm(0.77))
    text_frame_fourth = textbox_fourth.text_frame
    text_frame_fourth.word_wrap = True
    text_frame_fourth.auto_size = True
    p_fourth = text_frame_fourth.paragraphs[0]
    run_fourth = p_fourth.add_run()
    run_fourth.text = "*The indices in '%' are relative to the semester average." 
    font_fourth = run_fourth.font
    font_fourth.size = Pt(11)
    font_fourth.name = 'DM Sans'
    font_fourth.bold = False
    font_fourth.color.rgb = RGBColor(25, 46, 104)  # Cor do texto em preto


    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/4a - Card_VIEWS_v2.png'
    slide.shapes.add_picture(image_path, Cm(3.43), Cm(1.96), width=Cm(8.97), height=Cm(8.91))

    image_path = f'dados_full/{artista}/plots/4b - Card_REVENUE_v2.png'
    slide.shapes.add_picture(image_path, Cm(12.45), Cm(1.96), width=Cm(8.97), height=Cm(8.91))

    image_path = f'dados_full/{artista}/plots/4c - Card_SUBSCRIBERS_v2.png'
    slide.shapes.add_picture(image_path, Cm(21.42), Cm(1.96), width=Cm(8.97), height=Cm(8.91))

    image_path = f'dados_full/{artista}/plots/4d - Card_RPM_v2.png'
    slide.shapes.add_picture(image_path, Cm(3.43), Cm(10.08), width=Cm(8.97), height=Cm(8.91))

    image_path = f'dados_full/{artista}/plots/4e - Card_IMPRESSIONS_v2.png'
    slide.shapes.add_picture(image_path, Cm(12.45), Cm(10.08), width=Cm(8.97), height=Cm(8.91))

    image_path = f'dados_full/{artista}/plots/4f - Card_WATCHTIME_(HOURS)_v2.png'
    slide.shapes.add_picture(image_path, Cm(21.42), Cm(10.08), width=Cm(8.97), height=Cm(8.91))


    #SLIDE MÉTRICAS VOD AVANÇADAS ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Adicionar um novo slide para 'Receita'
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image7.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "General VOD Metrics"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(13.1), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/1 - Metrics_VOD_Advanced.png'
    slide.shapes.add_picture(image_path, Cm(0), Cm(3.16), width=Cm(33.92), height=Cm(13.64))


    #SLIDE MÉTRICAS LIVES AVANÇADAS ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Adicionar um novo slide para 'Receita'
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image7.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "General Live Metrics"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(13.1), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/2 - Metrics_Lives_Advanced.png'
    slide.shapes.add_picture(image_path, Cm(0), Cm(3.16), width=Cm(33.92), height=Cm(13.64))


    #SLIDE MÉTRICAS SHORTS AVANÇADAS ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Adicionar um novo slide para 'Receita'
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image7.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "General Shorts Metrics"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(13.1), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/3 - Metrics_Shorts_Advanced.png'
    slide.shapes.add_picture(image_path, Cm(0), Cm(3.16), width=Cm(33.92), height=Cm(13.64))


    # SLIDE ANALISE INICIAL ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Adicionar um novo slide para 'Análise Inicial'
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image7.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Initial Analysis"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(8.7), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/6 - Initial_Analysis.png'
    slide.shapes.add_picture(image_path, Cm(12.74), Cm(3.04), width=Cm(20.5), height=Cm(11.93))

    image_path_3 = f'dados_full/{artista}/plots/5 - Published.png'
    slide.shapes.add_picture(image_path_3, Cm(12.7), Cm(-0.36), width=Cm(19.49), height=Cm(4.44))

    image_path_2 = f'dados_full/{artista}/plots/7 - Watchtime.png'
    slide.shapes.add_picture(image_path_2, Cm(12.7), Cm(14.1), width=Cm(19.49), height=Cm(5.42))

    # Adicionar a terceira caixa de texto do novo slide
    textbox_fourth = slide.shapes.add_textbox(Cm(12.74), Cm(14.15), Cm(11), Cm(0.73))
    text_frame_fourth = textbox_fourth.text_frame
    text_frame_fourth.word_wrap = True
    text_frame_fourth.auto_size = True
    p_fourth = text_frame_fourth.paragraphs[0]
    run_fourth = p_fourth.add_run()
    run_fourth.text = "*Data refers only to lives and VODs" 
    font_fourth = run_fourth.font
    font_fourth.size = Pt(10)
    font_fourth.name = 'DM Sans'
    font_fourth.bold = False
    font_fourth.color.rgb = RGBColor(25, 46, 104)  # Cor do texto em preto


    # SLIDE MONETIZAÇÃO VELHO NOVO ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Adicionar um novo slide para 'Monetização'
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image15.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Revenue Old x New"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(12.23), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/8 - Monetization_v2.png'
    slide.shapes.add_picture(image_path, Cm(4.53), Cm(2.89), width=Cm(24.8), height=Cm(14.23))

    # Adicionar a terceira caixa de texto do novo slide
    textbox_fourth = slide.shapes.add_textbox(Cm(12.49), Cm(15.9), Cm(11), Cm(0.73))
    text_frame_fourth = textbox_fourth.text_frame
    text_frame_fourth.word_wrap = True
    text_frame_fourth.auto_size = True
    p_fourth = text_frame_fourth.paragraphs[0]
    run_fourth = p_fourth.add_run()
    run_fourth.text = "*Data refers only to lives and VODs" 
    font_fourth = run_fourth.font
    font_fourth.size = Pt(10)
    font_fourth.name = 'DM Sans'
    font_fourth.bold = False
    font_fourth.color.rgb = RGBColor(25, 46, 104)  # Cor do texto em preto


    # SLIDE Views VELHO NOVO ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Adicionar um novo slide para 'Views'
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image15.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Views Old x New"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(12.23), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/19 - Views_New_vs_Old.png'
    slide.shapes.add_picture(image_path, Cm(4.53), Cm(2.89), width=Cm(24.8), height=Cm(14.23))


    # SLIDE MONETIZAÇÃO POR FORMATOS ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Adicionar um novo slide para 'Monetização'
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image7.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Monetization by Formats"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(15.45), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/9 - Monetization by formats.png'
    slide.shapes.add_picture(image_path, Cm(4.53), Cm(2.89), width=Cm(24.8), height=Cm(14.23))
   
    # SLIDE CONVERSÃO ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Adicionar um novo slide para 'Conversão'
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image15.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Conversion"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(7.07), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/10 - Conversion.png'
    slide.shapes.add_picture(image_path, Cm(4.53), Cm(2.89), width=Cm(24.8), height=Cm(14.23))

    # Adicionar a terceira caixa de texto do novo slide
    textbox_fourth = slide.shapes.add_textbox(Cm(12.49), Cm(15.9), Cm(11), Cm(0.73))
    text_frame_fourth = textbox_fourth.text_frame
    text_frame_fourth.word_wrap = True
    text_frame_fourth.auto_size = True
    p_fourth = text_frame_fourth.paragraphs[0]
    run_fourth = p_fourth.add_run()
    run_fourth.text = "*Data refers only to lives and VODs" 
    font_fourth = run_fourth.font
    font_fourth.size = Pt(10)
    font_fourth.name = 'DM Sans'
    font_fourth.bold = False
    font_fourth.color.rgb = RGBColor(25, 46, 104)  # Cor do texto em preto


    # SLIDE QUALIDADE VOD ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Adicionar um novo slide para 'Qualidade'
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image7.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "VODs Quality"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(10.18), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))
    
    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/11 - Quality_Vod.png'
    slide.shapes.add_picture(image_path, Cm(0), Cm(4.72), width=Cm(17.06), height=Cm(9.61))

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/15 - Engagement_Vod.png'
    slide.shapes.add_picture(image_path, Cm(17.06), Cm(4.41), width=Cm(17.06), height=Cm(10.24))

    # Adicionar a terceira caixa de texto do novo slide
    textbox_fourth = slide.shapes.add_textbox(Cm(0), Cm(18.38), Cm(13.71), Cm(0.68))
    text_frame_fourth = textbox_fourth.text_frame
    text_frame_fourth.word_wrap = True
    text_frame_fourth.auto_size = True
    p_fourth = text_frame_fourth.paragraphs[0]
    run_fourth = p_fourth.add_run()
    run_fourth.text = "*Percentage data in parentheses is the average percentage watched" 
    font_fourth = run_fourth.font
    font_fourth.size = Pt(10)
    font_fourth.name = 'DM Sans'
    font_fourth.bold = False
    font_fourth.color.rgb = RGBColor(25, 46, 104)  # Cor do texto em preto


    # SLIDE QUALIDADE LIVES ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Adicionar um novo slide para 'Qualidade'
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image7.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Lives Quality"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(9.72), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/12 - Quality_Live.png'
    slide.shapes.add_picture(image_path, Cm(0), Cm(4.72), width=Cm(17.06), height=Cm(9.61))

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/16 - Engagement_Live.png'
    slide.shapes.add_picture(image_path, Cm(17.06), Cm(4.41), width=Cm(17.06), height=Cm(10.24))

    # Adicionar a terceira caixa de texto do novo slide
    textbox_fourth = slide.shapes.add_textbox(Cm(0), Cm(18.38), Cm(13.71), Cm(0.68))
    text_frame_fourth = textbox_fourth.text_frame
    text_frame_fourth.word_wrap = True
    text_frame_fourth.auto_size = True
    p_fourth = text_frame_fourth.paragraphs[0]
    run_fourth = p_fourth.add_run()
    run_fourth.text = "*Percentage data in parentheses is the average percentage watched" 
    font_fourth = run_fourth.font
    font_fourth.size = Pt(10)
    font_fourth.name = 'DM Sans'
    font_fourth.bold = False
    font_fourth.color.rgb = RGBColor(25, 46, 104)  # Cor do texto em preto


    # SLIDE QUALIDADE SHORTS ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Adicionar um novo slide para 'Qualidade'
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image7.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Shorts Quality"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(9.72), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/12.5 - Quality_Shorts.png'
    slide.shapes.add_picture(image_path, Cm(0), Cm(4.72), width=Cm(17.06), height=Cm(9.61))

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/16.5 - Engagement_Shorts.png'
    slide.shapes.add_picture(image_path, Cm(17.06), Cm(4.41), width=Cm(17.06), height=Cm(10.24))

    # Adicionar a terceira caixa de texto do novo slide
    textbox_fourth = slide.shapes.add_textbox(Cm(0), Cm(18.38), Cm(13.71), Cm(0.68))
    text_frame_fourth = textbox_fourth.text_frame
    text_frame_fourth.word_wrap = True
    text_frame_fourth.auto_size = True
    p_fourth = text_frame_fourth.paragraphs[0]
    run_fourth = p_fourth.add_run()
    run_fourth.text = "*Percentage data in parentheses is the average percentage watched" 
    font_fourth = run_fourth.font
    font_fourth.size = Pt(10)
    font_fourth.name = 'DM Sans'
    font_fourth.bold = False
    font_fourth.color.rgb = RGBColor(25, 46, 104)  # Cor do texto em preto


    # SLIDE COMUNIDADE ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Adicionar um novo slide para 'Comunidade'
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image7.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Community"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(9.72), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/17 - Community.png'
    slide.shapes.add_picture(image_path, Cm(4.53), Cm(2.89), width=Cm(24.8), height=Cm(14.23))


    #SLIDE ORIGEM DO TRÁFEGO ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Adicionar um novo slide para 'Receita'
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image15.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Traffic Source"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(11.06), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/13 - Traffic_Source.png'
    slide.shapes.add_picture(image_path, Cm(4.53), Cm(2.89), width=Cm(24.8), height=Cm(14.23))

    # Adicionar a terceira caixa de texto do novo slide
    textbox_fourth = slide.shapes.add_textbox(Cm(12.49), Cm(15.9), Cm(11), Cm(0.73))
    text_frame_fourth = textbox_fourth.text_frame
    text_frame_fourth.word_wrap = True
    text_frame_fourth.auto_size = True
    p_fourth = text_frame_fourth.paragraphs[0]
    run_fourth = p_fourth.add_run()
    run_fourth.text = "*Data refers only to lives and VODs" 
    font_fourth = run_fourth.font
    font_fourth.size = Pt(10)
    font_fourth.name = 'DM Sans'
    font_fourth.bold = False
    font_fourth.color.rgb = RGBColor(25, 46, 104)  # Cor do texto em preto

    #SLIDE CRESCIMENTO ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Adicionar um novo slide para 'Crescimento'
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image15.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Subscriber Growth"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(14.69), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/14 - Subscribers_by_Content_Type.png'
    slide.shapes.add_picture(image_path, Cm(6.23), Cm(1.96), width=Cm(21.4), height=Cm(12.28))

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/18 - Subscribers Table.png'
    slide.shapes.add_picture(image_path, Cm(7.76), Cm(13.28), width=Cm(18.35), height=Cm(6.56))

    #SLIDE OBRIGADO ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path_1 = 'resources/Imagens Template Relatório Mensal/image32.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path_1, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a imagem
    image_path = 'resources/Imagens Template Relatório Mensal/image3.png'
    slide.shapes.add_picture(image_path, Cm(1.42), Cm(1.21), width=Cm(2.89), height=Cm(2.06))
    
    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(0.98), Cm(7.8), Cm(26.39), Cm(4.53))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Thank you!"
    font_first = run_first.font
    font_first.size = Pt(50)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(255, 255, 255)  # Cor do texto em branco

    # Adicionar a segunda caixa de texto do novo slide
    textbox_second = slide.shapes.add_textbox(Cm(1.95), Cm(14.45), Cm(14.34), Cm(1.95))
    text_frame_second = textbox_second.text_frame
    text_frame_second.word_wrap = True
    p_second = text_frame_second.paragraphs[0]
    run_second = p_second.add_run()
    run_second.text = "@1bigmedia\n" + "facebook/onebigmedia\n" + "youtube.com/onebigmedia\n"
    font_second = run_second.font
    font_second.size = Pt(14)
    font_second.bold = True
    font_second.name = 'DM Sans'
    font_second.color.rgb = RGBColor(255, 255, 255)  # Cor do texto em branco
    p_second.line_spacing = Pt(26.6)

    insta_path = f'resources/Imagens Template Relatório Mensal/image34.png'
    slide.shapes.add_picture(insta_path, Cm(1.42), Cm(14.88), width=Cm(0.67), height=Cm(0.67))

    facebook_path = f'resources/Imagens Template Relatório Mensal/image28.png'
    slide.shapes.add_picture(facebook_path, Cm(1.42), Cm(15.79), width=Cm(0.67), height=Cm(0.67))

    youtube_path = f'resources/Imagens Template Relatório Mensal/image30.png'
    slide.shapes.add_picture(youtube_path, Cm(1.37), Cm(16.65), width=Cm(0.77), height=Cm(0.74))


    # Salvar a apresentação
    prs.save(f'export_teste/Report Mensal {mes_foco} {artist}.pptx')


Perfeito! Vamos corrigir o apresentacao_report_ingles.py seguindo o mesmo padrão.

Com base no trecho que você forneceu, aqui estão as alterações necessárias:

Adicione import sys e import json no início do arquivo apresentacao_report_ingles.py.

Remova a função buscar_lista_artistas() (você precisará localizá-la e apagá-la, assim como fez nos outros scripts).

Substitua a função run() e o bloco if __name__ == "__main__": existentes pelo código corrigido abaixo:

Python

import sys
import json # Necessário para ler comparacoes.json
# ... (outros imports do seu script, como pptx, os, etc.)

# ... (função create_apresentation e outras funções auxiliares, provavelmente com nomes/textos em inglês) ...

def run(artista_tupla):
    """
    Executes the generation of the 'report' presentation (English) for a single artist.
    Receives a tuple (artist_file_name, artist_display_name).
    """
    # Desempacota a tupla recebida
    artista_arquivo = artista_tupla[0]
    artista_display = artista_tupla[1] # 'artist' no seu código original

    # Define o mês ou outras variáveis
    mes_foco = "August 2025" # Mantenha ou ajuste (já estava em inglês)

    print(f'Generating report presentation (English) for: {artista_display}')

    try:
        # A chamada para a função principal que cria a apresentação (em inglês)
        # Use as variáveis desempacotadas corretamente
        create_apresentation(artista_arquivo, artista_display, mes_foco) # Certifique-se que esta função gere o conteúdo em inglês

        print(f'-> Success: Report presentation (English) for {artista_display} completed.') # Mensagem de sucesso

    except FileNotFoundError as e:
        print(f"ERROR: File not found while generating report presentation (English) for {artista_arquivo}. Details: {e}")
    except Exception as e:
        print(f"ERROR: Unexpected error generating report presentation (English) for {artista_arquivo}: {e}")


if __name__ == "__main__":
    # Verifica se o main.py passou o nome do artista (nome do arquivo)
    if len(sys.argv) < 2:
        print("Error: No artist (file_name) provided. Should be called via main.py.")
        sys.exit(1)

    artista_arquivo_arg = sys.argv[1]

    # Tenta encontrar o nome de display correspondente em comparacoes.json
    try:
        # Garanta que o encoding esteja correto (latin-1 ou utf-8)
        with open('comparacoes.json', encoding='latin-1') as json_file:
            comparacoes = json.load(json_file)
        # Usa o nome do arquivo como fallback se não encontrar
        artista_display_arg = comparacoes.get(artista_arquivo_arg, artista_arquivo_arg)
    except FileNotFoundError:
        print("Warning: 'comparacoes.json' not found. Using file name as display name.")
        artista_display_arg = artista_arquivo_arg
    except Exception as e:
        print(f"Error reading 'comparacoes.json': {e}. Using file name as display name.")
        artista_display_arg = artista_arquivo_arg

    # Cria a tupla (nome_arquivo, nome_display) que a função run espera
    artista_tupla_arg = (artista_arquivo_arg, artista_display_arg)

    # Executa a função run APENAS para o artista fornecido
    run(artista_tupla_arg)








