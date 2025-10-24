import os
import sys
import json
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Cm
from pptx.enum.shapes import MSO_CONNECTOR
import re  
import docx



def extrair_conteudo_markdown(caminho_do_arquivo):
    """
    Versão final e mais completa da função para extrair APENAS OS BLOCOS DE TEXTO.
    Suporta três padrões de título:
    1. ### 1. Título
    2. **1. Título**
    3. 1. Título

    Args:
      caminho_do_arquivo: O caminho para o arquivo .txt.

    Returns:
      Uma lista de strings, onde cada string é um bloco de texto de uma seção.
    """
    try:
        with open(caminho_do_arquivo, 'r', encoding='latin-1') as file:
            conteudo = file.read()
    except FileNotFoundError:
        print(f"Erro: Arquivo não encontrado no caminho: {caminho_do_arquivo}")
        return []
    except Exception as e:
        print(f"Ocorreu um erro inesperado: {e}")
        return []

    # Regex Final: Unifica os 3 padrões.
    # A principal mudança é tornar o prefixo (### ou **) opcional com (?:...)?
    padrao = re.compile(
        # --- Início da Linha do Título ---
        r"^(?:###\s+|\*\*\s*)?"           # Grupo Opcional: Encontra '### ' OU '** ' (ou nada).
        r"(\d+\.\s+.*?)"                  # Grupo 1: Captura o número e o texto do título (será descartado).
        r"(?:\*\*|$)"                     # Encontra o '**' final OU o fim da linha.
        # --- Fim da Linha do Título ---
        r"[\r\n]+"                        # Corresponde a uma ou mais quebras de linha.
        r"([\s\S]*?)"                     # Grupo 2: Captura o Bloco de Texto (o que queremos).
        # --- Lookahead para o Fim do Bloco ---
        r"(?=^(?:###\s+|\*\*\s*)?\d+\.|\Z)", # Para a captura quando encontrar o início de QUALQUER um 
                                          # dos 3 padrões de título ou o fim do arquivo.
        re.MULTILINE
    )

    matches = padrao.findall(conteudo)

    # Cria a lista final pegando apenas o segundo elemento (o texto) de cada tupla.
    textos_extraidos = [texto.strip() for titulo, texto in matches]
    
    return textos_extraidos


def create_apresentation(artista, artist, mes_foco, textos_extraidos, graficos_path):
    prs = Presentation()

    # SLIDE CAPA ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Definir o tamanho do slide para widescreen (16:9)
    prs.slide_width = Inches(13.333)  # 13.333in é o equivalente a 16:9 em widescreen
    prs.slide_height = Inches(7.5)  # Altura proporcional para o tamanho widescreen

    # Caminho para a imagem de fundo do slide de capa
    background_image_path = 'resources/Imagens Template Relatório Mensal/image35.png'

    # Adicionar um slide de capa
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Adicionar a imagem de fundo
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a imagem
    image_path = 'resources/Imagens Template Relatório Mensal/image3.png'
    slide.shapes.add_picture(image_path, Cm(1.42), Cm(1.21), width=Cm(2.89), height=Cm(2.06))

    # Adicionar a primeira caixa de texto (título)
    textbox_title = slide.shapes.add_textbox(Cm(1.04), Cm(11.28), Cm(25.2), Cm(3.68))
    text_frame_title = textbox_title.text_frame
    p_title = text_frame_title.paragraphs[0]
    run_title = p_title.add_run()
    run_title.text = "Monthly Report"

    # Configuração da fonte do título
    font_title = run_title.font
    font_title.size = Pt(56)
    font_title.bold = True
    font_title.name = 'DM Sans'
    font_title.color.rgb = RGBColor(255, 255, 255)  # Cor do título em formato RGB hexadecimal

    # Adicionar a segunda caixa de texto (subtítulo)
    textbox_subtitle = slide.shapes.add_textbox(Cm(1.04), Cm(13.75), Cm(13.31), Cm(1.28))
    text_frame_subtitle = textbox_subtitle.text_frame
    p_subtitle = text_frame_subtitle.paragraphs[0]
    run_subtitle = p_subtitle.add_run()
    run_subtitle.text = mes_foco

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(1.34), Cm(15.61), Cm(17), Cm(15.61))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(255, 255, 255)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Configuração da fonte do subtítulo
    font_subtitle = run_subtitle.font
    font_subtitle.size = Pt(24)
    font_subtitle.bold = True
    font_subtitle.name = 'DM Sans'
    font_subtitle.color.rgb = RGBColor(255, 255, 255)  # Cor do subtítulo em formato RGB hexadecimal (#ffffff)

    # Adicionar a caixa de artista
    textbox_artist = slide.shapes.add_textbox(Cm(1.04), Cm(16.57), Cm(31.98), Cm(1.54))
    text_frame_artist = textbox_artist.text_frame
    p_artista = text_frame_artist.paragraphs[0]
    run_artist = p_artista.add_run()
    run_artist.text = artist

    # Configuração da fonte do subtítulo
    font_artist = run_artist.font
    font_artist.size = Pt(30)
    font_artist.bold = True
    font_artist.name = 'DM Sans'
    font_artist.color.rgb = RGBColor(255, 255, 255)  # Cor do subtítulo em formato RGB hexadecimal (#ffffff)
   
        #SLIDE DIAGNÓSTICO ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path_1 = 'resources/Imagens Template Relatório Mensal/image27.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path_1, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Slide Editável"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(7.12), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    # Adicionar a terceira caixa de texto do novo slide
    textbox_third = slide.shapes.add_textbox(Cm(1.21), Cm(2.89), Cm(11.01), Cm(1.64))
    text_frame_third = textbox_third.text_frame
    text_frame_third.word_wrap = True
    text_frame_third.auto_size = True
    p_third = text_frame_third.paragraphs[0]
    run_third = p_third.add_run()
    run_third.text = "Texto" #arquivo[arquivo['Tópicos'] == 'Diagnóstico']['Textos'].iloc[0]
    font_third = run_third.font
    font_third.size = Pt(14)
    font_third.name = 'DM Sans'
    font_third.color.rgb = RGBColor(0, 0, 0)  # Cor do texto em preto
    
    
    # SLIDE RESUMO ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image15.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Analysis Summary"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(11.32), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))


    # Adicionar a terceira caixa de texto do novo slide
    textbox_fourth = slide.shapes.add_textbox(Cm(0.50), Cm(18.34), Cm(11), Cm(0.77))
    text_frame_fourth = textbox_fourth.text_frame
    text_frame_fourth.word_wrap = True
    text_frame_fourth.auto_size = True
    p_fourth = text_frame_fourth.paragraphs[0]
    run_fourth = p_fourth.add_run()
    run_fourth.text = "*The indices in '%' are relative to the semester average." 
    font_fourth = run_fourth.font
    font_fourth.size = Pt(11)
    font_fourth.name = 'DM Sans'
    font_fourth.bold = False
    font_fourth.color.rgb = RGBColor(25, 46, 104)  # Cor do texto em preto


    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/4a - Card_VIEWS_v2.png'
    slide.shapes.add_picture(image_path, Cm(3.43), Cm(1.96), width=Cm(8.97), height=Cm(8.91))

    image_path = f'dados_full/{artista}/plots/4b - Card_REVENUE_v2.png'
    slide.shapes.add_picture(image_path, Cm(12.45), Cm(1.96), width=Cm(8.97), height=Cm(8.91))

    image_path = f'dados_full/{artista}/plots/4c - Card_SUBSCRIBERS_v2.png'
    slide.shapes.add_picture(image_path, Cm(21.42), Cm(1.96), width=Cm(8.97), height=Cm(8.91))

    image_path = f'dados_full/{artista}/plots/4d - Card_RPM_v2.png'
    slide.shapes.add_picture(image_path, Cm(3.43), Cm(10.08), width=Cm(8.97), height=Cm(8.91))

    image_path = f'dados_full/{artista}/plots/4e - Card_IMPRESSIONS_v2.png'
    slide.shapes.add_picture(image_path, Cm(12.45), Cm(10.08), width=Cm(8.97), height=Cm(8.91))

    image_path = f'dados_full/{artista}/plots/4f - Card_WATCHTIME_(HOURS)_v2.png'
    slide.shapes.add_picture(image_path, Cm(21.42), Cm(10.08), width=Cm(8.97), height=Cm(8.91))


    #SLIDE MÉTRICAS VOD AVANÇADAS ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Adicionar um novo slide para 'Receita'
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image7.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "General VOD Metrics"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(13.1), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/1 - Metrics_VOD_Advanced.png'
    slide.shapes.add_picture(image_path, Cm(0), Cm(3.16), width=Cm(33.92), height=Cm(13.64))


    #SLIDE MÉTRICAS LIVES AVANÇADAS ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Adicionar um novo slide para 'Receita'
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image7.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "General Live Metrics"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(13.1), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/2 - Metrics_Lives_Advanced.png'
    slide.shapes.add_picture(image_path, Cm(0), Cm(3.16), width=Cm(33.92), height=Cm(13.64))


    #SLIDE MÉTRICAS SHORTS AVANÇADAS ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Adicionar um novo slide para 'Receita'
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image7.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "General Shorts Metrics"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(13.1), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/3 - Metrics_Shorts_Advanced.png'
    slide.shapes.add_picture(image_path, Cm(0), Cm(3.16), width=Cm(33.92), height=Cm(13.64))


    # SLIDE ANALISE INICIAL ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Adicionar um novo slide para 'Análise Inicial'
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image7.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Initial Analysis"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(8.7), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    # Adicionar a terceira caixa de texto do novo slide
    textbox_third = slide.shapes.add_textbox(Cm(1.21), Cm(2.89), Cm(11.01), Cm(1.64))
    text_frame_third = textbox_third.text_frame
    text_frame_third.word_wrap = True
    text_frame_third.auto_size = True
    p_third = text_frame_third.paragraphs[0]
    run_third = p_third.add_run()
    run_third.text = "Texto" #arquivo[arquivo['Tópicos'] == 'Análise Inicial']['Textos'].iloc[0]
    font_third = run_third.font
    font_third.size = Pt(14)
    font_third.name = 'DM Sans'
    font_third.color.rgb = RGBColor(0, 0, 0)  # Cor do texto em preto

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/6 - Initial_Analysis.png'
    slide.shapes.add_picture(image_path, Cm(12.74), Cm(3.04), width=Cm(20.5), height=Cm(11.93))

    image_path_3 = f'dados_full/{artista}/plots/5 - Published.png'
    slide.shapes.add_picture(image_path_3, Cm(12.7), Cm(-0.36), width=Cm(19.49), height=Cm(4.44))

    image_path_2 = f'dados_full/{artista}/plots/7 - Watchtime.png'
    slide.shapes.add_picture(image_path_2, Cm(12.7), Cm(14.1), width=Cm(19.49), height=Cm(5.42))

    # Adicionar a terceira caixa de texto do novo slide
    textbox_fourth = slide.shapes.add_textbox(Cm(12.74), Cm(14.15), Cm(11), Cm(0.73))
    text_frame_fourth = textbox_fourth.text_frame
    text_frame_fourth.word_wrap = True
    text_frame_fourth.auto_size = True
    p_fourth = text_frame_fourth.paragraphs[0]
    run_fourth = p_fourth.add_run()
    run_fourth.text = "*Data refers only to lives and VODs" 
    font_fourth = run_fourth.font
    font_fourth.size = Pt(10)
    font_fourth.name = 'DM Sans'
    font_fourth.bold = False
    font_fourth.color.rgb = RGBColor(25, 46, 104)  # Cor do texto em preto


    # SLIDE MONETIZAÇÃO VELHO NOVO ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Adicionar um novo slide para 'Monetização'
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image15.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Revenue Old x New"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(12.23), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    # Adicionar a terceira caixa de texto do novo slide
    textbox_third = slide.shapes.add_textbox(Cm(1.21), Cm(2.89), Cm(11.01), Cm(1.64))
    text_frame_third = textbox_third.text_frame
    text_frame_third.word_wrap = True
    text_frame_third.auto_size = True
    p_third = text_frame_third.paragraphs[0]
    run_third = p_third.add_run()
    run_third.text = "Texto" #arquivo[arquivo['Tópicos'] == 'Monetização Velho x Novo']['Textos'].iloc[0]
    font_third = run_third.font
    font_third.size = Pt(14)
    font_third.name = 'DM Sans'
    font_third.color.rgb = RGBColor(0, 0, 0)  # Cor do texto em preto

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/8 - Monetization_v2.png'
    slide.shapes.add_picture(image_path, Cm(12.49), Cm(3.62), width=Cm(21.4), height=Cm(12.28))

    # Adicionar a terceira caixa de texto do novo slide
    textbox_fourth = slide.shapes.add_textbox(Cm(12.49), Cm(15.9), Cm(11), Cm(0.73))
    text_frame_fourth = textbox_fourth.text_frame
    text_frame_fourth.word_wrap = True
    text_frame_fourth.auto_size = True
    p_fourth = text_frame_fourth.paragraphs[0]
    run_fourth = p_fourth.add_run()
    run_fourth.text = "*Data refers only to lives and VODs" 
    font_fourth = run_fourth.font
    font_fourth.size = Pt(10)
    font_fourth.name = 'DM Sans'
    font_fourth.bold = False
    font_fourth.color.rgb = RGBColor(25, 46, 104)  # Cor do texto em preto


    # SLIDE Views VELHO NOVO ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Adicionar um novo slide para 'Views'
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image15.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Views Old x New"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(12.23), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    # Adicionar a terceira caixa de texto do novo slide
    textbox_third = slide.shapes.add_textbox(Cm(1.21), Cm(2.89), Cm(11.01), Cm(1.64))
    text_frame_third = textbox_third.text_frame
    text_frame_third.word_wrap = True
    text_frame_third.auto_size = True
    p_third = text_frame_third.paragraphs[0]
    run_third = p_third.add_run()
    run_third.text = "Texto" #arquivo[arquivo['Tópicos'] == 'Monetização Velho x Novo']['Textos'].iloc[0]
    font_third = run_third.font
    font_third.size = Pt(14)
    font_third.name = 'DM Sans'
    font_third.color.rgb = RGBColor(0, 0, 0)  # Cor do texto em preto

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/19 - Views_New_vs_Old.png'
    slide.shapes.add_picture(image_path, Cm(12.49), Cm(3.62), width=Cm(21.4), height=Cm(12.28))


    # SLIDE MONETIZAÇÃO POR FORMATOS ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Adicionar um novo slide para 'Monetização'
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image7.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Monetization by Formats"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(15.45), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    # Adicionar a terceira caixa de texto do novo slide
    textbox_third = slide.shapes.add_textbox(Cm(1.21), Cm(2.89), Cm(11.01), Cm(1.64))
    text_frame_third = textbox_third.text_frame
    text_frame_third.word_wrap = True
    text_frame_third.auto_size = True
    p_third = text_frame_third.paragraphs[0]
    run_third = p_third.add_run()
    run_third.text = "Texto" #arquivo[arquivo['Tópicos'] == 'Monetização por Formatos']['Textos'].iloc[0]
    font_third = run_third.font
    font_third.size = Pt(14)
    font_third.name = 'DM Sans'
    font_third.color.rgb = RGBColor(0, 0, 0)  # Cor do texto em preto

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/9 - Monetization by formats.png'
    slide.shapes.add_picture(image_path, Cm(12.49), Cm(3.62), width=Cm(21.4), height=Cm(12.28))
   
    # SLIDE CONVERSÃO ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Adicionar um novo slide para 'Conversão'
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image15.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Conversion"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(7.07), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    # Adicionar a terceira caixa de texto do novo slide
    textbox_third = slide.shapes.add_textbox(Cm(1.21), Cm(2.89), Cm(11.01), Cm(1.64))
    text_frame_third = textbox_third.text_frame
    text_frame_third.word_wrap = True
    text_frame_third.auto_size = True
    p_third = text_frame_third.paragraphs[0]
    run_third = p_third.add_run()
    run_third.text = "Texto" #arquivo[arquivo['Tópicos'] == 'Conversão']['Textos'].iloc[0]
    font_third = run_third.font
    font_third.size = Pt(14)
    font_third.name = 'DM Sans'
    font_third.color.rgb = RGBColor(0, 0, 0)  # Cor do texto em preto

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/10 - Conversion.png'
    slide.shapes.add_picture(image_path, Cm(12.49), Cm(3.62), width=Cm(21.4), height=Cm(12.28))

    # Adicionar a terceira caixa de texto do novo slide
    textbox_fourth = slide.shapes.add_textbox(Cm(12.49), Cm(15.9), Cm(11), Cm(0.73))
    text_frame_fourth = textbox_fourth.text_frame
    text_frame_fourth.word_wrap = True
    text_frame_fourth.auto_size = True
    p_fourth = text_frame_fourth.paragraphs[0]
    run_fourth = p_fourth.add_run()
    run_fourth.text = "*Data refers only to lives and VODs" 
    font_fourth = run_fourth.font
    font_fourth.size = Pt(10)
    font_fourth.name = 'DM Sans'
    font_fourth.bold = False
    font_fourth.color.rgb = RGBColor(25, 46, 104)  # Cor do texto em preto


    # SLIDE QUALIDADE VOD ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Adicionar um novo slide para 'Qualidade'
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image7.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "VODs Quality"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(10.18), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    # Adicionar a terceira caixa de texto do novo slide
    textbox_third = slide.shapes.add_textbox(Cm(1.21), Cm(2.89), Cm(11.01), Cm(1.64))
    text_frame_third = textbox_third.text_frame
    text_frame_third.word_wrap = True
    text_frame_third.auto_size = True
    p_third = text_frame_third.paragraphs[0]
    run_third = p_third.add_run()
    run_third.text = "Texto" #arquivo[arquivo['Tópicos'] == "Qualidade VOD's"]['Textos'].iloc[0]
    font_third = run_third.font
    font_third.size = Pt(14)
    font_third.name = 'DM Sans'
    font_third.color.rgb = RGBColor(0, 0, 0)  # Cor do texto em preto
    
    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/11 - Quality_Vod.png'
    slide.shapes.add_picture(image_path, Cm(14.51), Cm(-0.11), width=Cm(17.06), height=Cm(9.61))

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/15 - Engagement_Vod.png'
    slide.shapes.add_picture(image_path, Cm(14.51), Cm(9.32), width=Cm(17.06), height=Cm(10.24))

    # Adicionar a terceira caixa de texto do novo slide
    textbox_fourth = slide.shapes.add_textbox(Cm(0), Cm(18.38), Cm(13.71), Cm(0.68))
    text_frame_fourth = textbox_fourth.text_frame
    text_frame_fourth.word_wrap = True
    text_frame_fourth.auto_size = True
    p_fourth = text_frame_fourth.paragraphs[0]
    run_fourth = p_fourth.add_run()
    run_fourth.text = "*Percentage data in parentheses is the average percentage watched" 
    font_fourth = run_fourth.font
    font_fourth.size = Pt(10)
    font_fourth.name = 'DM Sans'
    font_fourth.bold = False
    font_fourth.color.rgb = RGBColor(25, 46, 104)  # Cor do texto em preto

    # SLIDE QUALIDADE LIVES ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Adicionar um novo slide para 'Qualidade'
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image7.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Lives Quality"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(9.72), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    # Adicionar a terceira caixa de texto do novo slide
    textbox_third = slide.shapes.add_textbox(Cm(1.21), Cm(2.89), Cm(11.01), Cm(1.64))
    text_frame_third = textbox_third.text_frame
    text_frame_third.word_wrap = True
    text_frame_third.auto_size = True
    p_third = text_frame_third.paragraphs[0]
    run_third = p_third.add_run()
    run_third.text = "Texto" #arquivo[arquivo['Tópicos'] == 'Qualidade Lives']['Textos'].iloc[0]
    font_third = run_third.font
    font_third.size = Pt(14)
    font_third.name = 'DM Sans'
    font_third.color.rgb = RGBColor(0, 0, 0)  # Cor do texto em preto

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/12 - Quality_Live.png'
    slide.shapes.add_picture(image_path, Cm(14.51), Cm(-0.11), width=Cm(17.06), height=Cm(9.61))

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/16 - Engagement_Live.png'
    slide.shapes.add_picture(image_path, Cm(14.51), Cm(9.32), width=Cm(17.06), height=Cm(10.24))

    # Adicionar a terceira caixa de texto do novo slide
    textbox_fourth = slide.shapes.add_textbox(Cm(0), Cm(18.38), Cm(13.71), Cm(0.68))
    text_frame_fourth = textbox_fourth.text_frame
    text_frame_fourth.word_wrap = True
    text_frame_fourth.auto_size = True
    p_fourth = text_frame_fourth.paragraphs[0]
    run_fourth = p_fourth.add_run()
    run_fourth.text = "*Percentage data in parentheses is the average percentage watched" 
    font_fourth = run_fourth.font
    font_fourth.size = Pt(10)
    font_fourth.name = 'DM Sans'
    font_fourth.bold = False
    font_fourth.color.rgb = RGBColor(25, 46, 104)  # Cor do texto em preto


    # SLIDE QUALIDADE SHORTS ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Adicionar um novo slide para 'Qualidade'
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image7.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Shorts Quality"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(9.72), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    # Adicionar a terceira caixa de texto do novo slide
    textbox_third = slide.shapes.add_textbox(Cm(1.21), Cm(2.89), Cm(11.01), Cm(1.64))
    text_frame_third = textbox_third.text_frame
    text_frame_third.word_wrap = True
    text_frame_third.auto_size = True
    p_third = text_frame_third.paragraphs[0]
    run_third = p_third.add_run()
    run_third.text = "Texto" #arquivo[arquivo['Tópicos'] == 'Qualidade Lives']['Textos'].iloc[0]
    font_third = run_third.font
    font_third.size = Pt(14)
    font_third.name = 'DM Sans'
    font_third.color.rgb = RGBColor(0, 0, 0)  # Cor do texto em preto

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/12.5 - Quality_Shorts.png'
    slide.shapes.add_picture(image_path, Cm(14.51), Cm(-0.11), width=Cm(17.06), height=Cm(9.61))

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/16.5 - Engagement_Shorts.png'
    slide.shapes.add_picture(image_path, Cm(14.51), Cm(9.32), width=Cm(17.06), height=Cm(10.24))

    # Adicionar a terceira caixa de texto do novo slide
    textbox_fourth = slide.shapes.add_textbox(Cm(0), Cm(18.38), Cm(13.71), Cm(0.68))
    text_frame_fourth = textbox_fourth.text_frame
    text_frame_fourth.word_wrap = True
    text_frame_fourth.auto_size = True
    p_fourth = text_frame_fourth.paragraphs[0]
    run_fourth = p_fourth.add_run()
    run_fourth.text = "*Percentage data in parentheses is the average percentage watched" 
    font_fourth = run_fourth.font
    font_fourth.size = Pt(10)
    font_fourth.name = 'DM Sans'
    font_fourth.bold = False
    font_fourth.color.rgb = RGBColor(25, 46, 104)  # Cor do texto em preto

    # SLIDE COMUNIDADE ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Adicionar um novo slide para 'Comunidade'
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image7.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Community"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(9.72), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    # Adicionar a terceira caixa de texto do novo slide
    textbox_third = slide.shapes.add_textbox(Cm(1.21), Cm(2.89), Cm(11.01), Cm(1.64))
    text_frame_third = textbox_third.text_frame
    text_frame_third.word_wrap = True
    text_frame_third.auto_size = True
    p_third = text_frame_third.paragraphs[0]
    run_third = p_third.add_run()
    run_third.text = "Texto" #arquivo[arquivo['Tópicos'] == 'Qualidade Lives']['Textos'].iloc[0]
    font_third = run_third.font
    font_third.size = Pt(14)
    font_third.name = 'DM Sans'
    font_third.color.rgb = RGBColor(0, 0, 0)  # Cor do texto em preto

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/17 - Community.png'
    slide.shapes.add_picture(image_path, Cm(12.49), Cm(3.62), width=Cm(21.4), height=Cm(12.28))


    #SLIDE ORIGEM DO TRÁFEGO ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Adicionar um novo slide para 'Receita'
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image15.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Traffic Source"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(11.06), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    # Adicionar a terceira caixa de texto do novo slide
    textbox_third = slide.shapes.add_textbox(Cm(1.21), Cm(2.89), Cm(11.01), Cm(1.64))
    text_frame_third = textbox_third.text_frame
    text_frame_third.word_wrap = True
    text_frame_third.auto_size = True
    p_third = text_frame_third.paragraphs[0]
    run_third = p_third.add_run()
    run_third.text = "Texto" #arquivo[arquivo['Tópicos'] == 'Monetização Velho x Novo']['Textos'].iloc[0]
    font_third = run_third.font
    font_third.size = Pt(14)
    font_third.name = 'DM Sans'
    font_third.color.rgb = RGBColor(0, 0, 0)  # Cor do texto em preto

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/13 - Traffic_Source.png'
    slide.shapes.add_picture(image_path, Cm(12.49), Cm(3.62), width=Cm(21.4), height=Cm(12.28))

    # Adicionar a terceira caixa de texto do novo slide
    textbox_fourth = slide.shapes.add_textbox(Cm(12.49), Cm(15.9), Cm(11), Cm(0.73))
    text_frame_fourth = textbox_fourth.text_frame
    text_frame_fourth.word_wrap = True
    text_frame_fourth.auto_size = True
    p_fourth = text_frame_fourth.paragraphs[0]
    run_fourth = p_fourth.add_run()
    run_fourth.text = "*Data refers only to lives and VODs" 
    font_fourth = run_fourth.font
    font_fourth.size = Pt(10)
    font_fourth.name = 'DM Sans'
    font_fourth.bold = False
    font_fourth.color.rgb = RGBColor(25, 46, 104)  # Cor do texto em preto

    #SLIDE CRESCIMENTO ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    # Adicionar um novo slide para 'Crescimento'
    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path = 'resources/Imagens Template Relatório Mensal/image15.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(18), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Subscriber Growth"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(14.69), Cm(2.1))

    # Formatar a linha (opcional)
    line_shape.line.color.rgb = RGBColor(49, 87, 247)  # Cor da linha em azul escuro
    line_shape.line.width = Pt(3)  # Espessura da linha

    # Adicionar a imagem
    logo_path = f'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    # Adicionar a terceira caixa de texto do novo slide
    textbox_third = slide.shapes.add_textbox(Cm(1.21), Cm(2.89), Cm(11.01), Cm(1.64))
    text_frame_third = textbox_third.text_frame
    text_frame_third.word_wrap = True
    text_frame_third.auto_size = True
    p_third = text_frame_third.paragraphs[0]
    run_third = p_third.add_run()
    run_third.text = "Texto" #arquivo[arquivo['Tópicos'] == 'Crescimento']['Textos'].iloc[0]
    font_third = run_third.font
    font_third.size = Pt(14)
    font_third.name = 'DM Sans'
    font_third.color.rgb = RGBColor(0, 0, 0)  # Cor do texto em preto

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/14 - Subscribers_by_Content_Type.png'
    slide.shapes.add_picture(image_path, Cm(12.49), Cm(1.52), width=Cm(21.4), height=Cm(12.28))

    # Adicionar a imagem
    image_path = f'dados_full/{artista}/plots/18 - Subscribers Table.png'
    slide.shapes.add_picture(image_path, Cm(14.01), Cm(12.96), width=Cm(18.35), height=Cm(6.56))

    ###################################################################################
    # SLIDE Cluster - Quantidade de Conteúdos #
    ###################################################################################
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    background_image_path = 'resources/Imagens Template Relatório Mensal/image15.png'
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(25), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Cluster - Content Quantity"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(20.23), Cm(2.1))
    line_shape.line.color.rgb = RGBColor(49, 87, 247)
    line_shape.line.width = Pt(3)

    logo_path = 'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    textbox_third = slide.shapes.add_textbox(Cm(1.21), Cm(2.89), Cm(11), Cm(15)) # Aumentei o tamanho da caixa de texto
    text_frame_third = textbox_third.text_frame
    text_frame_third.word_wrap = True
    p_third = text_frame_third.paragraphs[0]
    run_third = p_third.add_run()
    
    # [NOVO] Pega o texto do dicionário. Usa get() para evitar erros se a chave não existir.
    run_third.text = textos_extraidos[0]
    
    font_third = run_third.font
    font_third.size = Pt(14)
    font_third.name = 'DM Sans'
    font_third.color.rgb = RGBColor(0, 0, 0)

    # Adicionar a imagem
    image_path = f'{graficos_path}grafico_publicacoes_{artista}.png'
    slide.shapes.add_picture(image_path, Cm(12.49), Cm(3.62), width=Cm(21.4), height=Cm(12.28))

    
    ###################################################################################
    # SLIDE Cluster - Visualizações #
    ###################################################################################
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    background_image_path = 'resources/Imagens Template Relatório Mensal/image15.png'
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(25), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Cluster - Views"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(12.23), Cm(2.1))
    line_shape.line.color.rgb = RGBColor(49, 87, 247)
    line_shape.line.width = Pt(3)

    logo_path = 'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    textbox_third = slide.shapes.add_textbox(Cm(1.21), Cm(2.89), Cm(11), Cm(15))
    text_frame_third = textbox_third.text_frame
    text_frame_third.word_wrap = True
    p_third = text_frame_third.paragraphs[0]
    run_third = p_third.add_run()

    # [NOVO] Pega o texto do dicionário.
    run_third.text = textos_extraidos[1]

    font_third = run_third.font
    font_third.size = Pt(14)
    font_third.name = 'DM Sans'
    font_third.color.rgb = RGBColor(0, 0, 0)

    # Adicionar a imagem
    image_path = f'{graficos_path}grafico_visualizacoes_{artista}.png'
    slide.shapes.add_picture(image_path, Cm(12.49), Cm(3.62), width=Cm(21.4), height=Cm(12.28))


    ###################################################################################
    # SLIDE Cluster - Receita Estimada #
    ###################################################################################
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    background_image_path = 'resources/Imagens Template Relatório Mensal/image15.png'
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(25), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Cluster - Estimated Revenue"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(14.23), Cm(2.1))
    line_shape.line.color.rgb = RGBColor(49, 87, 247)
    line_shape.line.width = Pt(3)

    logo_path = 'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    textbox_third = slide.shapes.add_textbox(Cm(1.21), Cm(2.89), Cm(11), Cm(15))
    text_frame_third = textbox_third.text_frame
    text_frame_third.word_wrap = True
    p_third = text_frame_third.paragraphs[0]
    run_third = p_third.add_run()
    
    # [NOVO] Pega o texto do dicionário.
    run_third.text = textos_extraidos[2]

    font_third = run_third.font
    font_third.size = Pt(14)
    font_third.name = 'DM Sans'
    font_third.color.rgb = RGBColor(0, 0, 0)

    # Adicionar a imagem
    image_path = f'{graficos_path}grafico_receita_{artista}.png'
    slide.shapes.add_picture(image_path, Cm(12.49), Cm(3.62), width=Cm(21.4), height=Cm(12.28))


    ###################################################################################
    # SLIDE Cluster - Média de Visualização #
    ###################################################################################
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    background_image_path = 'resources/Imagens Template Relatório Mensal/image15.png'
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(25), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Cluster - Average View"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(18.23), Cm(2.1))
    line_shape.line.color.rgb = RGBColor(49, 87, 247)
    line_shape.line.width = Pt(3)

    logo_path = 'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    textbox_third = slide.shapes.add_textbox(Cm(1.21), Cm(2.89), Cm(11), Cm(15))
    text_frame_third = textbox_third.text_frame
    text_frame_third.word_wrap = True
    p_third = text_frame_third.paragraphs[0]
    run_third = p_third.add_run()
    
    # [NOVO] Pega o texto do dicionário.
    run_third.text = textos_extraidos[3]

    font_third = run_third.font
    font_third.size = Pt(14)
    font_third.name = 'DM Sans'
    font_third.color.rgb = RGBColor(0, 0, 0)

    # Adicionar a imagem
    image_path = f'{graficos_path}grafico_media_visualizacoes_{artista}.png'
    slide.shapes.add_picture(image_path, Cm(12.49), Cm(3.62), width=Cm(21.4), height=Cm(12.28))


    ###################################################################################
    # SLIDE Cluster - Média de Receita #
    ###################################################################################
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    background_image_path = 'resources/Imagens Template Relatório Mensal/image15.png'
    slide.shapes.add_picture(background_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    textbox_first = slide.shapes.add_textbox(Cm(1.21), Cm(0.42), Cm(25), Cm(1.54))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Cluster - Average Revenue"
    font_first = run_first.font
    font_first.size = Pt(30)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)

    line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(-0.01), Cm(2.1), Cm(14.23), Cm(2.1))
    line_shape.line.color.rgb = RGBColor(49, 87, 247)
    line_shape.line.width = Pt(3)

    logo_path = 'resources/Imagens Template Relatório Mensal/image1.png'
    slide.shapes.add_picture(logo_path, Cm(0.53), Cm(0.77), width=Cm(0.63), height=Cm(0.83))

    textbox_third = slide.shapes.add_textbox(Cm(1.21), Cm(2.89), Cm(11), Cm(15))
    text_frame_third = textbox_third.text_frame
    text_frame_third.word_wrap = True
    p_third = text_frame_third.paragraphs[0]
    run_third = p_third.add_run()
    
    # [NOVO] Pega o texto do dicionário.
    run_third.text = textos_extraidos[4]
    
    font_third = run_third.font
    font_third.size = Pt(14)
    font_third.name = 'DM Sans'
    font_third.color.rgb = RGBColor(0, 0, 0)

    # Adicionar a imagem
    image_path = f'{graficos_path}grafico_media_receita_{artista}.png'
    slide.shapes.add_picture(image_path, Cm(12.49), Cm(3.62), width=Cm(21.4), height=Cm(12.28))


    #SLIDE DUVIDAS ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path_1 = 'resources/Imagens Template Relatório Mensal/image31.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path_1, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(8.54), Cm(7.26), Cm(16.83), Cm(4.53))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Questions?"
    font_first = run_first.font
    font_first.size = Pt(100)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(49, 87, 247)  # Cor do texto em azul escuro

    #SLIDE OBRIGADO ------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------

    slide_layout = prs.slide_layouts[5]  # Escolhendo um layout de slide em branco
    slide = prs.slides.add_slide(slide_layout)

    # Caminho para a imagem de fundo do novo slide
    background_image_path_1 = 'resources/Imagens Template Relatório Mensal/image32.png'

    # Adicionar a imagem de fundo do novo slide
    slide.shapes.add_picture(background_image_path_1, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Adicionar a imagem
    image_path = 'resources/Imagens Template Relatório Mensal/image3.png'
    slide.shapes.add_picture(image_path, Cm(1.42), Cm(1.21), width=Cm(2.89), height=Cm(2.06))
    
    # Adicionar a primeira caixa de texto do novo slide
    textbox_first = slide.shapes.add_textbox(Cm(0.98), Cm(7.8), Cm(26.39), Cm(4.53))
    text_frame_first = textbox_first.text_frame
    p_first = text_frame_first.paragraphs[0]
    run_first = p_first.add_run()
    run_first.text = "Thank you!"
    font_first = run_first.font
    font_first.size = Pt(50)
    font_first.bold = True
    font_first.name = 'DM Sans'
    font_first.color.rgb = RGBColor(255, 255, 255)  # Cor do texto em branco

    # Adicionar a segunda caixa de texto do novo slide
    textbox_second = slide.shapes.add_textbox(Cm(1.95), Cm(14.45), Cm(14.34), Cm(1.95))
    text_frame_second = textbox_second.text_frame
    text_frame_second.word_wrap = True
    p_second = text_frame_second.paragraphs[0]
    run_second = p_second.add_run()
    run_second.text = "@1bigmedia\n" + "facebook/onebigmedia\n" + "youtube.com/onebigmedia\n"
    font_second = run_second.font
    font_second.size = Pt(14)
    font_second.bold = True
    font_second.name = 'DM Sans'
    font_second.color.rgb = RGBColor(255, 255, 255)  # Cor do texto em branco
    p_second.line_spacing = Pt(26.6)

    insta_path = f'resources/Imagens Template Relatório Mensal/image34.png'
    slide.shapes.add_picture(insta_path, Cm(1.42), Cm(14.88), width=Cm(0.67), height=Cm(0.67))

    facebook_path = f'resources/Imagens Template Relatório Mensal/image28.png'
    slide.shapes.add_picture(facebook_path, Cm(1.42), Cm(15.79), width=Cm(0.67), height=Cm(0.67))

    youtube_path = f'resources/Imagens Template Relatório Mensal/image30.png'
    slide.shapes.add_picture(youtube_path, Cm(1.37), Cm(16.65), width=Cm(0.77), height=Cm(0.74))


    # Salvar a apresentação
    prs.save(f'export_teste/Overview Mensal {mes_foco} {artist}.pptx')


# (Certifique-se de que 'import sys' está no topo)
# (Certifique-se de que 'buscar_lista_artistas()' foi removida)

def run(artista_para_processar_tupla): # <-- MUDANÇA: Recebe a tupla
    # O loop 'for' foi REMOVIDO
    
    # Desempacota a tupla
    artista_arquivo = artista_para_processar_tupla[0] 
    artista_display = artista_para_processar_tupla[1]

    # Ajuste o mês ou torne-o dinâmico se necessário
    mes_foco = "September 2025" # Exemplo em inglês

    # === Verifique e ajuste estes caminhos ===
    # Caminho para o arquivo de texto (pode ser o mesmo ou uma versão em inglês)
    caminho_texto = rf'C:\Users\samuj\OneDrive\Área de Trabalho\1bmg\projeto_cluster_API\relatorios_api/{artista_arquivo}/exports/analise_narrativa_{artista_arquivo}_en.txt' # Exemplo com sufixo _en
    # Caminho para os gráficos (pode ser o mesmo ou uma pasta diferente)
    graficos_path = rf'C:\Users\samuj\OneDrive\Área de Trabalho\1bmg\projeto_cluster_API\relatorios_api/{artista_arquivo}/exports/' # Ajuste se necessário
    # === Fim da verificação de caminhos ===

    print(f"Generating 'cluster' presentation (English) for: {artista_display}")
    
    try:
        # Tenta extrair o conteúdo do markdown (ajuste a função se necessário)
        # Assumindo que a função extrair_conteudo_markdown funciona para inglês também
        textos_extraidos = extrair_conteudo_markdown(caminho_texto) 
        if not textos_extraidos:
             print(f"WARNING: Could not extract text from markdown for {artista_arquivo}. Check file: {caminho_texto}")
             # Use placeholders or decide to stop
             textos_extraidos = ["Text not found"] * 5 

        # Chama a função principal para criar a apresentação
        # Assumindo que existe uma 'create_presentation_ingles' ou que a original serve
        create_presentation_ingles(artista_arquivo, artista_display, mes_foco, textos_extraidos, graficos_path) # Ajuste o nome da função se necessário
        
        print(f"-> Success: 'cluster' presentation (English) for {artista_display} generated.")

    except FileNotFoundError as e:
        print(f"ERROR: File not found while generating presentation for {artista_arquivo}. Details: {e}")
    except Exception as e:
        print(f"ERROR: Unexpected error generating 'cluster' presentation (English) for {artista_arquivo}: {e}")


if __name__ == "__main__":
    
    if len(sys.argv) < 2:
        print("Error: No artist provided (expected format: file_name). This script must be called by main.py")
        sys.exit(1) 
    
    artista_arquivo_arg = sys.argv[1]
    try:
        with open('comparacoes.json', encoding='latin-1') as json_file: # Use o encoding correto
            comparacoes = json.load(json_file)
        # Busca o nome de display; se não encontrar, usa o nome do arquivo
        artista_display_arg = comparacoes.get(artista_arquivo_arg, artista_arquivo_arg) 
    except FileNotFoundError:
        print("Warning: 'comparacoes.json' not found. Using file name as display name.")
        artista_display_arg = artista_arquivo_arg
    except Exception as e:
        print(f"Error reading 'comparacoes.json': {e}. Using file name as display name.")
        artista_display_arg = artista_arquivo_arg

    # Cria a tupla esperada pela função run
    artista_tupla_arg = (artista_arquivo_arg, artista_display_arg)
    
    # Executa a função run APENAS para esse artista
    run(artista_tupla_arg)
